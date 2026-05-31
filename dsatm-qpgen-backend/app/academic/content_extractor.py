#!/usr/bin/env python3
"""
=============================================================================
  Content Extractor — Standalone Document Extraction & Image-Keyword Mapper
=============================================================================

Extracts text and images from PDFs, Word documents, PowerPoint files,
and standalone images. Maps each extracted image to related keywords
using RAKE + TF-IDF, with optional AI vision captioning (BLIP-2).

Usage:
    python content_extractor.py <input_file> [options]

Examples:
    python content_extractor.py textbook.pdf
    python content_extractor.py notes.docx --output-dir ./output
    python content_extractor.py diagram.png --ocr-lang eng+hin
    python content_extractor.py slides.pptx --use-ai

Output:
    - extracted_content.json  (structured extraction result)
    - extracted_images/       (all images saved as PNG files)

Author:  QPGen Content Extraction Module
Version: 2.0.0  (smart filtering)
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import io
import json
import logging
import os
import re
import sys
import time
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

# ---------------------------------------------------------------------------
#  Third-party imports (with graceful fallback messages)
# ---------------------------------------------------------------------------

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document as DocxDocument
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

try:
    from PIL import Image, ImageEnhance, ImageFilter
except ImportError:
    Image = None

try:
    import pytesseract
    # Windows fallback: point to default installation path if it exists
    if sys.platform == "win32":
        default_tess_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(default_tess_path):
            pytesseract.pytesseract.tesseract_cmd = default_tess_path
except ImportError:
    pytesseract = None

try:
    from rake_nltk import Rake
except ImportError:
    Rake = None

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
except ImportError:
    TfidfVectorizer = None

# ---------------------------------------------------------------------------
#  Logging setup
# ---------------------------------------------------------------------------

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger("ContentExtractor")


# ===========================================================================
#  Data Models
# ===========================================================================

@dataclass
class TextBlock:
    """Represents a block of extracted text."""
    page: int
    text: str
    block_index: int
    source_type: str = ""  # "paragraph", "table", "slide", "ocr"


@dataclass
class ExtractedImage:
    """Represents an extracted image with keyword mapping."""
    image_id: str
    image_path: str
    source_page: int
    width: int = 0
    height: int = 0
    related_text: str = ""
    keywords: list[str] = field(default_factory=list)
    ai_caption: str = ""
    context_before: str = ""
    context_after: str = ""
    image_hash: str = ""


@dataclass
class ExtractionResult:
    """Complete extraction result for a document."""
    source_file: str
    file_type: str
    total_pages: int = 0
    extraction_timestamp: str = ""
    text_blocks: list[TextBlock] = field(default_factory=list)
    images: list[ExtractedImage] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)

    def to_dict(self) -> dict:
        return {
            "source_file": self.source_file,
            "file_type": self.file_type,
            "total_pages": self.total_pages,
            "extraction_timestamp": self.extraction_timestamp,
            "text_blocks": [asdict(tb) for tb in self.text_blocks],
            "images": [asdict(img) for img in self.images],
            "metadata": self.metadata,
        }

    def to_json(self, indent: int = 2) -> str:
        return json.dumps(self.to_dict(), indent=indent, ensure_ascii=False)


# ===========================================================================
#  Content Filter (Smart Filtering)
# ===========================================================================

class ContentFilter:
    """
    Intelligently filters extracted content to keep only material
    useful for question paper generation. Removes:
      - Page headers/footers (repeated across pages)
      - Page numbers
      - Section/chapter titles (short standalone headings)
      - Figure/table captions ("Figure 1.1 ...", "Table 2 ...")
      - Table of contents entries
      - Copyright notices and institutional boilerplate
      - Very short fragments (< 30 chars) that are labels/bullets
      - Tiny images (inline formulas, icons, bullets)
      - Duplicate images (same hash)
    """

    # Minimum image dimensions to be considered a real diagram/figure
    MIN_IMAGE_WIDTH = 200
    MIN_IMAGE_HEIGHT = 150
    MIN_IMAGE_AREA = 40000  # width * height

    # Minimum text length to keep (shorter blocks are usually labels)
    MIN_TEXT_LENGTH = 30

    # Patterns that indicate non-content text (headers, titles, captions, etc.)
    NOISE_PATTERNS = [
        # Page numbers
        re.compile(r'^\s*\d{1,3}\s*$'),
        re.compile(r'^\s*page\s+\d+', re.IGNORECASE),
        re.compile(r'^\s*-\s*\d+\s*-\s*$'),
        # Figure/table captions (keep these separate — they reference images)
        re.compile(r'^\s*(?:figure|fig\.?)\s+\d', re.IGNORECASE),
        re.compile(r'^\s*(?:table|tab\.?)\s+\d', re.IGNORECASE),
        # Section numbers only ("1.", "2.3", "iv.")
        re.compile(r'^\s*(?:[ivxlc]+\.?|\d{1,3}\.\d{0,3}\.?)\s*$', re.IGNORECASE),
        # "Module X" / "Chapter X" / "Unit X" standalone titles
        re.compile(r'^\s*(?:module|chapter|unit|section|part)\s+\d', re.IGNORECASE),
        # Copyright / institutional
        re.compile(r'(?:copyright|all rights reserved|\u00a9)', re.IGNORECASE),
        # "Dept of" / university headers
        re.compile(r'^\s*(?:dept\.?|department)\s+of', re.IGNORECASE),
    ]

    # Patterns for section headings (short title-like text)
    HEADING_PATTERNS = [
        # Numbered headings: "1. What is AI?", "ii. Thinking rationally"
        re.compile(r'^\s*(?:[ivxlc]+|\d{1,3})\.\s+\S', re.IGNORECASE),
    ]

    def __init__(self, enabled: bool = True):
        self.enabled = enabled
        self._page_text_counts: dict[str, int] = {}  # text -> count across pages

    def filter_result(self, result: "ExtractionResult") -> "ExtractionResult":
        """Apply all smart filters to an extraction result."""
        if not self.enabled:
            return result

        original_text_count = len(result.text_blocks)
        original_image_count = len(result.images)

        # Step 1: Detect repeated headers/footers across pages
        self._detect_repeated_text(result.text_blocks)

        # Step 2: Filter text blocks
        result.text_blocks = self._filter_text_blocks(result.text_blocks)

        # Step 3: Filter images (size + dedup)
        result.images = self._filter_images(result.images)

        filtered_text = original_text_count - len(result.text_blocks)
        filtered_images = original_image_count - len(result.images)
        logger.info(
            f"Smart filter: removed {filtered_text} text blocks, "
            f"{filtered_images} images (kept {len(result.text_blocks)} text, "
            f"{len(result.images)} images)"
        )

        return result

    def _detect_repeated_text(self, text_blocks: list["TextBlock"]):
        """Find text that appears on multiple pages (likely headers/footers)."""
        self._page_text_counts.clear()
        page_texts: dict[str, set[int]] = {}  # normalized_text -> set of pages

        for tb in text_blocks:
            # Normalize: lowercase, strip, collapse whitespace
            normalized = re.sub(r'\s+', ' ', tb.text.strip().lower())
            if len(normalized) < 5:
                continue
            if normalized not in page_texts:
                page_texts[normalized] = set()
            page_texts[normalized].add(tb.page)

        # Text appearing on 3+ pages (or >40% of pages) is likely a header/footer
        total_pages = len(set(tb.page for tb in text_blocks))
        threshold = max(3, int(total_pages * 0.4))

        for text, pages in page_texts.items():
            if len(pages) >= threshold:
                self._page_text_counts[text] = len(pages)

        if self._page_text_counts:
            logger.debug(
                f"Detected {len(self._page_text_counts)} repeated header/footer patterns"
            )

    def _filter_text_blocks(self, text_blocks: list["TextBlock"]) -> list["TextBlock"]:
        """Filter out non-content text blocks."""
        filtered = []

        for tb in text_blocks:
            text = tb.text.strip()

            # --- Skip very short text ---
            if len(text) < self.MIN_TEXT_LENGTH:
                logger.debug(f"Filtered (too short, {len(text)} chars): {text[:50]}")
                continue

            # --- Skip repeated headers/footers ---
            normalized = re.sub(r'\s+', ' ', text.lower())
            if normalized in self._page_text_counts:
                logger.debug(f"Filtered (repeated header/footer): {text[:50]}")
                continue

            # --- Skip noise patterns ---
            is_noise = False
            for pattern in self.NOISE_PATTERNS:
                if pattern.search(text):
                    is_noise = True
                    logger.debug(f"Filtered (noise pattern): {text[:50]}")
                    break
            if is_noise:
                continue

            # --- Skip standalone short headings (< 80 chars, looks like a title) ---
            if len(text) < 80 and self._is_heading(text):
                logger.debug(f"Filtered (heading): {text[:50]}")
                continue

            filtered.append(tb)

        return filtered

    def _is_heading(self, text: str) -> bool:
        """Check if a short text block is likely a heading/title."""
        text = text.strip()

        # Numbered section heading: "1. What is AI?"
        for pattern in self.HEADING_PATTERNS:
            if pattern.match(text):
                return True

        # All-caps or title-case short text
        words = text.split()
        if len(words) <= 6:
            if text.isupper():
                return True
            # Title case: most words capitalized
            cap_count = sum(1 for w in words if w[0].isupper())
            if cap_count >= len(words) * 0.7 and not any(c in text for c in '.;:,'):
                return True

        return False

    def _filter_images(self, images: list["ExtractedImage"]) -> list["ExtractedImage"]:
        """Filter out tiny/junk images and deduplicate."""
        filtered = []
        seen_hashes: set[str] = set()

        for img in images:
            # --- Skip tiny images (inline formulas, icons, bullets) ---
            if img.width < self.MIN_IMAGE_WIDTH or img.height < self.MIN_IMAGE_HEIGHT:
                logger.debug(
                    f"Filtered image {img.image_id}: too small "
                    f"({img.width}x{img.height} < {self.MIN_IMAGE_WIDTH}x{self.MIN_IMAGE_HEIGHT})"
                )
                # Delete the file
                self._safe_delete(img.image_path)
                continue

            if img.width * img.height < self.MIN_IMAGE_AREA:
                logger.debug(
                    f"Filtered image {img.image_id}: area too small "
                    f"({img.width * img.height} < {self.MIN_IMAGE_AREA})"
                )
                self._safe_delete(img.image_path)
                continue

            # --- Skip duplicates ---
            if img.image_hash and img.image_hash in seen_hashes:
                logger.debug(f"Filtered image {img.image_id}: duplicate hash {img.image_hash}")
                self._safe_delete(img.image_path)
                continue

            if img.image_hash:
                seen_hashes.add(img.image_hash)

            filtered.append(img)

        # Re-number image IDs sequentially
        for i, img in enumerate(filtered, start=1):
            img.image_id = f"img_{i:04d}"

        return filtered



    def _safe_delete(self, path: str):
        """Safely delete a file."""
        try:
            if os.path.exists(path):
                os.remove(path)
        except Exception:
            pass


# ===========================================================================
#  Image-Keyword Mapper
# ===========================================================================

class ImageKeywordMapper:
    """
    Maps images to keywords extracted from surrounding text context.
    Uses RAKE (primary) + TF-IDF (secondary) for keyword extraction.
    """

    # Common stopwords to filter out from keywords
    EXTRA_STOPWORDS = {
        "figure", "fig", "image", "diagram", "table", "chart",
        "shown", "shows", "show", "see", "refer", "following",
        "above", "below", "left", "right", "page", "example",
        "given", "following", "also", "one", "two", "three",
        "using", "used", "use", "can", "may", "would", "could",
    }

    def __init__(self, max_keywords: int = 15):
        self.max_keywords = max_keywords
        self._rake = None
        self._tfidf = None

        if Rake is not None:
            try:
                self._rake = Rake(
                    min_length=1,
                    max_length=3,
                    include_repeated_phrases=False,
                )
            except Exception:
                # RAKE may need nltk stopwords; download them
                try:
                    import nltk
                    nltk.download("stopwords", quiet=True)
                    nltk.download("punkt_tab", quiet=True)
                    self._rake = Rake(
                        min_length=1,
                        max_length=3,
                        include_repeated_phrases=False,
                    )
                except Exception as e:
                    logger.warning(f"RAKE init failed: {e}. Using TF-IDF only.")

        if TfidfVectorizer is not None:
            self._tfidf = TfidfVectorizer(
                max_features=50,
                stop_words="english",
                ngram_range=(1, 2),
            )

    def extract_keywords(self, text: str) -> list[str]:
        """Extract keywords from text using RAKE + TF-IDF."""
        if not text or len(text.strip()) < 10:
            return []

        keywords: list[str] = []

        # --- RAKE extraction ---
        if self._rake is not None:
            try:
                self._rake.extract_keywords_from_text(text)
                ranked = self._rake.get_ranked_phrases_with_scores()
                for score, phrase in ranked[:self.max_keywords]:
                    cleaned = self._clean_keyword(phrase)
                    if cleaned and cleaned not in keywords:
                        keywords.append(cleaned)
            except Exception as e:
                logger.debug(f"RAKE extraction failed: {e}")

        # --- TF-IDF extraction (supplement) ---
        if self._tfidf is not None and len(keywords) < self.max_keywords:
            try:
                sentences = [s.strip() for s in text.split(".") if len(s.strip()) > 10]
                if len(sentences) >= 2:
                    matrix = self._tfidf.fit_transform(sentences)
                    feature_names = self._tfidf.get_feature_names_out()
                    scores = matrix.sum(axis=0).A1
                    ranked_indices = scores.argsort()[::-1]
                    for idx in ranked_indices[:self.max_keywords]:
                        kw = str(feature_names[idx]).strip()
                        cleaned = self._clean_keyword(kw)
                        if cleaned and cleaned not in keywords:
                            keywords.append(cleaned)
            except Exception as e:
                logger.debug(f"TF-IDF extraction failed: {e}")

        # --- Fallback: simple frequency-based ---
        if not keywords:
            keywords = self._simple_keyword_extraction(text)

        return keywords[:self.max_keywords]

    def _clean_keyword(self, keyword: str) -> str:
        """Clean and validate a keyword."""
        kw = keyword.strip().lower()
        # Remove pure numbers and very short keywords
        if len(kw) < 2 or kw.isdigit():
            return ""
        # Remove stopwords
        words = kw.split()
        words = [w for w in words if w not in self.EXTRA_STOPWORDS and len(w) > 1]
        if not words:
            return ""
        return " ".join(words)

    def _simple_keyword_extraction(self, text: str) -> list[str]:
        """Fallback keyword extraction using word frequency."""
        words = re.findall(r"\b[a-zA-Z]{3,}\b", text.lower())
        all_stop = self.EXTRA_STOPWORDS | {
            "the", "and", "for", "are", "but", "not", "you", "all",
            "her", "was", "its", "has", "had", "how", "may", "who",
            "this", "that", "with", "from", "have", "been", "will",
            "they", "each", "which", "their", "there", "what", "were",
            "when", "where", "than", "then", "them", "these", "those",
            "into", "over", "some", "such", "only", "other", "very",
        }
        words = [w for w in words if w not in all_stop]
        freq: dict[str, int] = {}
        for w in words:
            freq[w] = freq.get(w, 0) + 1
        sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
        return [w for w, _ in sorted_words[:self.max_keywords]]

    def map_image_to_context(
        self,
        text_blocks: list[TextBlock],
        image_page: int,
        block_index: int = -1,
        context_window: int = 2,
    ) -> tuple[str, str, str, list[str]]:
        """
        Given an image's page and optional block index, find the related text,
        context before/after, and extract keywords.

        Returns:
            (related_text, context_before, context_after, keywords)
        """
        # Collect text blocks from the same page
        same_page_blocks = [tb for tb in text_blocks if tb.page == image_page]
        if not same_page_blocks:
            # Expand to nearby pages
            same_page_blocks = [
                tb for tb in text_blocks
                if abs(tb.page - image_page) <= 1
            ]

        if not same_page_blocks:
            return "", "", "", []

        # Build the full text from this page
        page_text = " ".join(tb.text for tb in same_page_blocks)

        # Find context before and after
        all_blocks_sorted = sorted(text_blocks, key=lambda b: (b.page, b.block_index))

        context_before_parts = []
        context_after_parts = []

        if block_index >= 0:
            # We know where the image is relative to text blocks
            for tb in all_blocks_sorted:
                if tb.page == image_page and tb.block_index < block_index:
                    context_before_parts.append(tb.text)
                elif tb.page == image_page and tb.block_index > block_index:
                    context_after_parts.append(tb.text)
        else:
            # Use first/second half of page text as before/after
            mid = len(same_page_blocks) // 2
            context_before_parts = [tb.text for tb in same_page_blocks[:mid]]
            context_after_parts = [tb.text for tb in same_page_blocks[mid:]]

        context_before = " ".join(context_before_parts[-context_window:])
        context_after = " ".join(context_after_parts[:context_window])

        # Truncate contexts to reasonable length
        context_before = context_before[-500:] if len(context_before) > 500 else context_before
        context_after = context_after[:500] if len(context_after) > 500 else context_after

        # Extract keywords from the full page text
        keywords = self.extract_keywords(page_text)

        return page_text[:1000], context_before, context_after, keywords


# ===========================================================================
#  AI Vision Captioner (Optional — requires transformers + torch)
# ===========================================================================

class AIVisionCaptioner:
    """
    Optional AI-powered image captioning using BLIP-2.
    Disabled by default. Enable with --use-ai flag.
    Requires: transformers, torch, accelerate (pip install transformers torch accelerate)
    """

    def __init__(self):
        self.model = None
        self.processor = None
        self.available = False

    def load_model(self) -> bool:
        """Load the BLIP-2 model. Returns True if successful."""
        try:
            from transformers import BlipProcessor, BlipForConditionalGeneration
            import torch

            logger.info("Loading BLIP vision model (this may take a minute)...")
            model_name = "Salesforce/blip-image-captioning-base"  # ~1GB, lighter version
            self.processor = BlipProcessor.from_pretrained(model_name)
            self.model = BlipForConditionalGeneration.from_pretrained(model_name)

            # Use GPU if available
            if torch.cuda.is_available():
                self.model = self.model.to("cuda")
                logger.info("BLIP model loaded on GPU.")
            else:
                logger.info("BLIP model loaded on CPU.")

            self.available = True
            return True
        except ImportError:
            logger.warning(
                "AI captioning requires: pip install transformers torch\n"
                "Falling back to keyword-only mapping."
            )
            return False
        except Exception as e:
            logger.warning(f"Failed to load BLIP model: {e}")
            return False

    def caption_image(self, image_path: str) -> str:
        """Generate a caption for an image file."""
        if not self.available or self.model is None:
            return ""
        try:
            import torch
            img = Image.open(image_path).convert("RGB")
            inputs = self.processor(img, return_tensors="pt")
            if torch.cuda.is_available():
                inputs = {k: v.to("cuda") for k, v in inputs.items()}
            with torch.no_grad():
                output = self.model.generate(**inputs, max_new_tokens=50)
            caption = self.processor.decode(output[0], skip_special_tokens=True)
            return caption.strip()
        except Exception as e:
            logger.debug(f"Caption generation failed for {image_path}: {e}")
            return ""

    def caption_image_bytes(self, image_bytes: bytes) -> str:
        """Generate a caption from raw image bytes."""
        if not self.available or self.model is None:
            return ""
        try:
            import torch
            img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            inputs = self.processor(img, return_tensors="pt")
            if torch.cuda.is_available():
                inputs = {k: v.to("cuda") for k, v in inputs.items()}
            with torch.no_grad():
                output = self.model.generate(**inputs, max_new_tokens=50)
            caption = self.processor.decode(output[0], skip_special_tokens=True)
            return caption.strip()
        except Exception as e:
            logger.debug(f"Caption generation failed: {e}")
            return ""


# ===========================================================================
#  PDF Extractor
# ===========================================================================

class PDFExtractor:
    """
    Extracts text and images from PDF files using PyMuPDF (fitz).
    Handles both digital PDFs and scanned/image-based PDFs (via OCR fallback).
    """

    def __init__(self, output_dir: str, ocr_lang: str = "eng"):
        self.output_dir = output_dir
        self.ocr_lang = ocr_lang
        self.images_dir = os.path.join(output_dir, "extracted_images")
        os.makedirs(self.images_dir, exist_ok=True)

    def extract(self, file_path: str) -> ExtractionResult:
        """Extract all content from a PDF file."""
        if fitz is None:
            raise ImportError(
                "PyMuPDF is required for PDF extraction.\n"
                "Install it: pip install PyMuPDF"
            )

        logger.info(f"Extracting PDF: {file_path}")
        result = ExtractionResult(
            source_file=os.path.basename(file_path),
            file_type="pdf",
            extraction_timestamp=datetime.now().isoformat(),
        )

        doc = fitz.open(file_path)
        result.total_pages = len(doc)
        result.metadata = {
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", ""),
            "creator": doc.metadata.get("creator", ""),
            "page_count": len(doc),
        }

        image_counter = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_idx = page_num + 1  # 1-indexed

            # --- Extract text blocks ---
            text_dict = page.get_text("dict")
            block_index = 0

            for block in text_dict.get("blocks", []):
                if block.get("type") == 0:  # Text block
                    lines = block.get("lines", [])
                    block_text = ""
                    for line in lines:
                        spans = line.get("spans", [])
                        line_text = " ".join(span.get("text", "") for span in spans)
                        block_text += line_text + " "

                    block_text = block_text.strip()
                    if block_text and len(block_text) > 2:
                        result.text_blocks.append(TextBlock(
                            page=page_idx,
                            text=block_text,
                            block_index=block_index,
                            source_type="paragraph",
                        ))
                        block_index += 1

            # --- Check if page is scanned (very little text) ---
            page_text = page.get_text("text").strip()
            if len(page_text) < 20:
                # Likely a scanned page — try OCR
                ocr_text = self._ocr_page(page)
                if ocr_text and len(ocr_text) > 10:
                    result.text_blocks.append(TextBlock(
                        page=page_idx,
                        text=ocr_text,
                        block_index=block_index,
                        source_type="ocr",
                    ))
                    block_index += 1

            # --- Extract images ---
            image_list = page.get_images(full=True)
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if base_image is None:
                        continue

                    image_bytes = base_image["image"]
                    image_ext = base_image.get("ext", "png")
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)

                    # Skip obviously invalid images (< 30px)
                    # Real filtering happens in ContentFilter after extraction
                    if width < 30 or height < 30:
                        continue

                    # Generate image ID and save
                    image_counter += 1
                    img_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                    image_id = f"img_{image_counter:04d}"
                    image_filename = f"{image_id}_{img_hash}.png"
                    image_path = os.path.join(self.images_dir, image_filename)

                    # Convert to PNG for consistency
                    try:
                        pil_img = Image.open(io.BytesIO(image_bytes))
                        pil_img.save(image_path, "PNG")
                    except Exception:
                        # Fallback: save raw bytes
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)

                    result.images.append(ExtractedImage(
                        image_id=image_id,
                        image_path=image_path,
                        source_page=page_idx,
                        width=width,
                        height=height,
                        image_hash=img_hash,
                    ))

                except Exception as e:
                    logger.warning(f"Failed to extract image xref={xref} on page {page_idx}: {e}")

        doc.close()
        logger.info(
            f"PDF extraction complete: {len(result.text_blocks)} text blocks, "
            f"{len(result.images)} images from {result.total_pages} pages"
        )
        return result

    def _ocr_page(self, page) -> str:
        """Run OCR on a PDF page rendered as an image."""
        if pytesseract is None or Image is None:
            return ""
        try:
            # Render page to a high-res image
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            pil_img = Image.open(io.BytesIO(img_bytes))

            # Pre-process for OCR
            pil_img = pil_img.convert("L")  # Grayscale
            enhancer = ImageEnhance.Contrast(pil_img)
            pil_img = enhancer.enhance(1.5)

            text = pytesseract.image_to_string(pil_img, lang=self.ocr_lang)
            return text.strip()
        except Exception as e:
            logger.debug(f"OCR failed on page: {e}")
            return ""


# ===========================================================================
#  Word Document Extractor
# ===========================================================================

class DocxExtractor:
    """
    Extracts text and images from .docx files using python-docx.
    Maps images to surrounding paragraph text.
    """

    def __init__(self, output_dir: str):
        self.output_dir = output_dir
        self.images_dir = os.path.join(output_dir, "extracted_images")
        os.makedirs(self.images_dir, exist_ok=True)

    def extract(self, file_path: str) -> ExtractionResult:
        """Extract all content from a Word document."""
        if DocxDocument is None:
            raise ImportError(
                "python-docx is required for Word document extraction.\n"
                "Install it: pip install python-docx"
            )

        logger.info(f"Extracting DOCX: {file_path}")
        result = ExtractionResult(
            source_file=os.path.basename(file_path),
            file_type="docx",
            extraction_timestamp=datetime.now().isoformat(),
        )

        doc = DocxDocument(file_path)

        # --- Extract document properties ---
        try:
            props = doc.core_properties
            result.metadata = {
                "title": props.title or "",
                "author": props.author or "",
                "subject": props.subject or "",
                "created": str(props.created) if props.created else "",
                "modified": str(props.modified) if props.modified else "",
            }
        except Exception:
            result.metadata = {}

        # --- Extract paragraphs ---
        block_index = 0
        paragraph_map: dict[int, str] = {}  # block_index -> text

        for para in doc.paragraphs:
            text = para.text.strip()
            if text and len(text) > 1:
                result.text_blocks.append(TextBlock(
                    page=1,  # DOCX doesn't have pages; use 1
                    text=text,
                    block_index=block_index,
                    source_type="paragraph",
                ))
                paragraph_map[block_index] = text
                block_index += 1

        # --- Extract table content ---
        for table in doc.tables:
            table_text_parts = []
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    table_text_parts.append(" | ".join(row_cells))

            if table_text_parts:
                table_text = "\n".join(table_text_parts)
                result.text_blocks.append(TextBlock(
                    page=1,
                    text=table_text,
                    block_index=block_index,
                    source_type="table",
                ))
                block_index += 1

        # --- Extract images from document relationships ---
        image_counter = 0
        image_rels = {}

        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    content_type = image_part.content_type

                    # Skip tiny images
                    try:
                        pil_img = Image.open(io.BytesIO(image_bytes))
                        w, h = pil_img.size
                        if w < 50 or h < 50:
                            continue
                    except Exception:
                        w, h = 0, 0
                        pil_img = None

                    image_counter += 1
                    img_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                    image_id = f"img_{image_counter:04d}"
                    image_filename = f"{image_id}_{img_hash}.png"
                    image_path = os.path.join(self.images_dir, image_filename)

                    # Save as PNG
                    if pil_img is not None:
                        pil_img.save(image_path, "PNG")
                    else:
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)

                    result.images.append(ExtractedImage(
                        image_id=image_id,
                        image_path=image_path,
                        source_page=1,
                        width=w,
                        height=h,
                        image_hash=img_hash,
                    ))

                except Exception as e:
                    logger.warning(f"Failed to extract DOCX image: {e}")

        # --- Map images to nearest paragraphs ---
        # In DOCX, images are inline in paragraphs. Walk the XML to find them.
        self._map_images_to_paragraphs(doc, result)

        result.total_pages = 1  # DOCX is page-less
        logger.info(
            f"DOCX extraction complete: {len(result.text_blocks)} text blocks, "
            f"{len(result.images)} images"
        )
        return result

    def _map_images_to_paragraphs(self, doc, result: ExtractionResult):
        """Walk paragraphs to find which ones contain inline images and record block_index."""
        try:
            from docx.oxml.ns import qn
            img_block_map: dict[int, int] = {}  # image_index -> block_index
            img_rel_counter = 0
            block_idx = 0

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    block_idx += 1

                # Check for inline images in this paragraph
                inline_elements = para._element.findall(
                    f".//{qn('wp:inline')}"
                ) + para._element.findall(
                    f".//{qn('wp:anchor')}"
                )

                for _ in inline_elements:
                    if img_rel_counter < len(result.images):
                        # Update the image's block reference for context mapping
                        result.images[img_rel_counter].source_page = block_idx
                        img_rel_counter += 1

        except Exception as e:
            logger.debug(f"Paragraph-image mapping encountered an issue: {e}")


# ===========================================================================
#  PowerPoint Extractor
# ===========================================================================

class PPTXExtractor:
    """
    Extracts text and images from .pptx files using python-pptx.
    Maps images to slide text content.
    """

    def __init__(self, output_dir: str):
        self.output_dir = output_dir
        self.images_dir = os.path.join(output_dir, "extracted_images")
        os.makedirs(self.images_dir, exist_ok=True)

    def extract(self, file_path: str) -> ExtractionResult:
        """Extract all content from a PowerPoint file."""
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PowerPoint extraction.\n"
                "Install it: pip install python-pptx"
            )

        logger.info(f"Extracting PPTX: {file_path}")
        result = ExtractionResult(
            source_file=os.path.basename(file_path),
            file_type="pptx",
            extraction_timestamp=datetime.now().isoformat(),
        )

        prs = Presentation(file_path)
        result.total_pages = len(prs.slides)

        image_counter = 0
        seen_hashes: set[str] = set()

        for slide_num, slide in enumerate(prs.slides, start=1):
            block_index = 0

            # --- Extract text from shapes ---
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_parts = []
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            text_parts.append(para_text)

                    if text_parts:
                        full_text = "\n".join(text_parts)
                        result.text_blocks.append(TextBlock(
                            page=slide_num,
                            text=full_text,
                            block_index=block_index,
                            source_type="slide",
                        ))
                        block_index += 1

                # --- Extract tables ---
                if shape.has_table:
                    table = shape.table
                    table_text_parts = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            table_text_parts.append(" | ".join(cells))
                    if table_text_parts:
                        result.text_blocks.append(TextBlock(
                            page=slide_num,
                            text="\n".join(table_text_parts),
                            block_index=block_index,
                            source_type="table",
                        ))
                        block_index += 1

                # --- Extract images ---
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image_part = shape.image
                        image_bytes = image_part.blob
                        img_hash = hashlib.md5(image_bytes).hexdigest()[:8]

                        # Skip duplicates
                        if img_hash in seen_hashes:
                            continue
                        seen_hashes.add(img_hash)

                        # Check size
                        try:
                            pil_img = Image.open(io.BytesIO(image_bytes))
                            w, h = pil_img.size
                            if w < 50 or h < 50:
                                continue
                        except Exception:
                            w, h = 0, 0
                            pil_img = None

                        image_counter += 1
                        image_id = f"img_{image_counter:04d}"
                        image_filename = f"{image_id}_{img_hash}.png"
                        image_path = os.path.join(self.images_dir, image_filename)

                        if pil_img is not None:
                            pil_img.save(image_path, "PNG")
                        else:
                            with open(image_path, "wb") as f:
                                f.write(image_bytes)

                        result.images.append(ExtractedImage(
                            image_id=image_id,
                            image_path=image_path,
                            source_page=slide_num,
                            width=w,
                            height=h,
                            image_hash=img_hash,
                        ))

                    except Exception as e:
                        logger.warning(f"Failed to extract PPTX image on slide {slide_num}: {e}")

        logger.info(
            f"PPTX extraction complete: {len(result.text_blocks)} text blocks, "
            f"{len(result.images)} images from {result.total_pages} slides"
        )
        return result


# ===========================================================================
#  Standalone Image Extractor (OCR)
# ===========================================================================

class ImageExtractor:
    """
    Extracts text from standalone image files using Tesseract OCR.
    Supports: PNG, JPG, JPEG, TIFF, BMP, WEBP
    """

    SUPPORTED_FORMATS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"}

    def __init__(self, output_dir: str, ocr_lang: str = "eng"):
        self.output_dir = output_dir
        self.ocr_lang = ocr_lang
        self.images_dir = os.path.join(output_dir, "extracted_images")
        os.makedirs(self.images_dir, exist_ok=True)

    def extract(self, file_path: str) -> ExtractionResult:
        """Extract text from a standalone image file via OCR."""
        if Image is None:
            raise ImportError(
                "Pillow is required for image processing.\n"
                "Install it: pip install Pillow"
            )

        logger.info(f"Extracting Image (OCR): {file_path}")
        result = ExtractionResult(
            source_file=os.path.basename(file_path),
            file_type="image",
            total_pages=1,
            extraction_timestamp=datetime.now().isoformat(),
        )

        # Open and pre-process image
        pil_img = Image.open(file_path)
        w, h = pil_img.size

        result.metadata = {
            "width": w,
            "height": h,
            "format": pil_img.format or "unknown",
            "mode": pil_img.mode,
        }

        # Copy original image to output
        image_id = "img_0001"
        img_hash = hashlib.md5(open(file_path, "rb").read()).hexdigest()[:8]
        image_filename = f"{image_id}_{img_hash}.png"
        image_path = os.path.join(self.images_dir, image_filename)
        pil_img.save(image_path, "PNG")

        # --- Run OCR ---
        ocr_text = ""
        if pytesseract is not None:
            try:
                # Pre-process for better OCR
                processed = self._preprocess_for_ocr(pil_img)
                ocr_text = pytesseract.image_to_string(processed, lang=self.ocr_lang)
                ocr_text = ocr_text.strip()
            except Exception as e:
                logger.warning(f"OCR failed: {e}")
                logger.info(
                    "Make sure Tesseract is installed:\n"
                    "  Windows: https://github.com/UB-Mannheim/tesseract/wiki\n"
                    "  Linux: sudo apt install tesseract-ocr\n"
                    "  macOS: brew install tesseract"
                )
        else:
            logger.warning(
                "pytesseract not installed. Install it: pip install pytesseract\n"
                "Also install Tesseract OCR engine on your system."
            )

        if ocr_text:
            result.text_blocks.append(TextBlock(
                page=1,
                text=ocr_text,
                block_index=0,
                source_type="ocr",
            ))

        result.images.append(ExtractedImage(
            image_id=image_id,
            image_path=image_path,
            source_page=1,
            width=w,
            height=h,
            related_text=ocr_text[:1000] if ocr_text else "",
            image_hash=img_hash,
        ))

        logger.info(
            f"Image extraction complete: "
            f"{'OCR text extracted' if ocr_text else 'No OCR text'} "
            f"({len(ocr_text)} chars)"
        )
        return result

    def _preprocess_for_ocr(self, img: "Image.Image") -> "Image.Image":
        """Pre-process image for better OCR accuracy."""
        # Convert to grayscale
        processed = img.convert("L")

        # Enhance contrast
        enhancer = ImageEnhance.Contrast(processed)
        processed = enhancer.enhance(1.8)

        # Sharpen
        enhancer = ImageEnhance.Sharpness(processed)
        processed = enhancer.enhance(2.0)

        # Upscale small images
        w, h = processed.size
        if w < 1000 or h < 1000:
            scale = max(1000 / w, 1000 / h, 1.0)
            if scale > 1.0:
                new_w = int(w * scale)
                new_h = int(h * scale)
                processed = processed.resize(
                    (new_w, new_h), Image.LANCZOS
                )

        return processed


# ===========================================================================
#  Main Orchestrator
# ===========================================================================

class ContentExtractor:
    """
    Main orchestrator that detects file type, routes to the appropriate
    extractor, and runs keyword mapping + optional AI captioning.
    """

    EXTENSION_MAP = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",  # Will warn about .doc needing conversion
        ".pptx": "pptx",
        ".ppt": "pptx",
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".tiff": "image",
        ".tif": "image",
        ".bmp": "image",
        ".webp": "image",
    }

    def __init__(
        self,
        output_dir: str = "./extraction_output",
        ocr_lang: str = "eng",
        use_ai: bool = False,
        max_keywords: int = 15,
        smart_filter: bool = True,
    ):
        self.output_dir = os.path.abspath(output_dir)
        self.ocr_lang = ocr_lang
        self.use_ai = use_ai
        self.max_keywords = max_keywords
        self.smart_filter = smart_filter

        os.makedirs(self.output_dir, exist_ok=True)

        # Initialize components
        self.keyword_mapper = ImageKeywordMapper(max_keywords=max_keywords)
        self.content_filter = ContentFilter(enabled=smart_filter)
        self.ai_captioner = AIVisionCaptioner() if use_ai else None

        # Load AI model if requested
        if use_ai and self.ai_captioner is not None:
            self.ai_captioner.load_model()

    def extract(self, file_path: str) -> ExtractionResult:
        """
        Extract content from any supported file.
        Returns an ExtractionResult with text blocks, images, and keyword mappings.
        """
        file_path = os.path.abspath(file_path)

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()
        file_type = self.EXTENSION_MAP.get(ext)

        if file_type is None:
            raise ValueError(
                f"Unsupported file type: {ext}\n"
                f"Supported: {', '.join(sorted(self.EXTENSION_MAP.keys()))}"
            )

        if ext == ".doc":
            raise ValueError(
                "Legacy .doc format is not supported directly.\n"
                "Please convert to .docx first (open in Word and Save As .docx)."
            )

        if ext == ".ppt":
            raise ValueError(
                "Legacy .ppt format is not supported directly.\n"
                "Please convert to .pptx first (open in PowerPoint and Save As .pptx)."
            )

        logger.info(f"{'=' * 60}")
        logger.info(f"Starting extraction: {os.path.basename(file_path)}")
        logger.info(f"File type: {file_type} | Output: {self.output_dir}")
        logger.info(f"{'=' * 60}")

        start_time = time.time()

        # --- Route to appropriate extractor ---
        if file_type == "pdf":
            extractor = PDFExtractor(self.output_dir, self.ocr_lang)
        elif file_type == "docx":
            extractor = DocxExtractor(self.output_dir)
        elif file_type == "pptx":
            extractor = PPTXExtractor(self.output_dir)
        elif file_type == "image":
            extractor = ImageExtractor(self.output_dir, self.ocr_lang)
        else:
            raise ValueError(f"No extractor available for: {file_type}")

        result = extractor.extract(file_path)

        # --- Apply smart content filter ---
        if self.smart_filter:
            logger.info("Applying smart content filter...")
            result = self.content_filter.filter_result(result)

        # --- Run keyword mapping on all images ---
        logger.info("Mapping images to keywords...")
        for img in result.images:
            related_text, ctx_before, ctx_after, keywords = (
                self.keyword_mapper.map_image_to_context(
                    result.text_blocks,
                    img.source_page,
                )
            )
            img.related_text = related_text or img.related_text
            img.context_before = ctx_before
            img.context_after = ctx_after
            img.keywords = keywords

            # --- Optional AI captioning ---
            if (
                self.use_ai
                and self.ai_captioner is not None
                and self.ai_captioner.available
            ):
                caption = self.ai_captioner.caption_image(img.image_path)
                if caption:
                    img.ai_caption = caption
                    # Add caption words to keywords
                    caption_words = re.findall(r"\b[a-zA-Z]{3,}\b", caption.lower())
                    for word in caption_words:
                        if word not in img.keywords and len(img.keywords) < self.max_keywords:
                            img.keywords.append(word)

        elapsed = time.time() - start_time

        # --- Save results ---
        output_json_path = os.path.join(self.output_dir, "extracted_content.json")
        with open(output_json_path, "w", encoding="utf-8") as f:
            f.write(result.to_json(indent=2))

        logger.info(f"{'=' * 60}")
        logger.info(f"Extraction complete in {elapsed:.2f}s")
        logger.info(f"  Text blocks: {len(result.text_blocks)}")
        logger.info(f"  Images:      {len(result.images)}")
        logger.info(f"  Output JSON: {output_json_path}")
        logger.info(f"  Images dir:  {os.path.join(self.output_dir, 'extracted_images')}")
        logger.info(f"{'=' * 60}")

        return result

    def check_dependencies(self) -> dict[str, bool]:
        """Check which dependencies are available."""
        deps = {
            "PyMuPDF (fitz)": fitz is not None,
            "python-docx": DocxDocument is not None,
            "python-pptx": Presentation is not None,
            "Pillow": Image is not None,
            "pytesseract": pytesseract is not None,
            "rake-nltk": Rake is not None,
            "scikit-learn (TF-IDF)": TfidfVectorizer is not None,
        }
        return deps


# ===========================================================================
#  Pretty Printer (for terminal output)
# ===========================================================================

def print_result_summary(result: ExtractionResult):
    """Print a formatted summary of extraction results to the terminal."""
    print()
    print("=" * 60)
    print("          CONTENT EXTRACTION RESULTS                    ")
    print("=" * 60)
    print()
    print(f"  Source:     {result.source_file}")
    print(f"  Type:       {result.file_type.upper()}")
    print(f"  Pages:      {result.total_pages}")
    print(f"  Timestamp:  {result.extraction_timestamp}")
    print()

    # --- Text Blocks Summary ---
    print(f"  --- Text Blocks ({len(result.text_blocks)}) ------------------------")
    for i, tb in enumerate(result.text_blocks[:10]):  # Show first 10
        preview = tb.text[:80].replace("\n", " ")
        if len(tb.text) > 80:
            preview += "..."
        print(f"  [{tb.page:>3}] {preview}")
    if len(result.text_blocks) > 10:
        print(f"  ... and {len(result.text_blocks) - 10} more blocks")
    print()

    # --- Images Summary ---
    print(f"  --- Images ({len(result.images)}) -----------------------------")
    for img in result.images:
        print(f"  {img.image_id}  |  page {img.source_page}  |  {img.width}x{img.height}px")
        if img.keywords:
            kw_display = ", ".join(img.keywords[:8])
            if len(img.keywords) > 8:
                kw_display += f", +{len(img.keywords) - 8} more"
            print(f"           |  keywords: {kw_display}")
        if img.ai_caption:
            print(f"           |  caption:  {img.ai_caption}")
        print(f"           |  path:     {img.image_path}")
        print()

    print("=" * 60)


# ===========================================================================
#  CLI Entry Point
# ===========================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Content Extractor — Extract text & images with keyword mapping",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python content_extractor.py textbook.pdf
  python content_extractor.py notes.docx --output-dir ./output
  python content_extractor.py diagram.png --ocr-lang eng+hin
  python content_extractor.py slides.pptx --use-ai
  python content_extractor.py physics.pdf --no-filter   # disable smart filtering
  python content_extractor.py physics.pdf --max-keywords 20 --verbose
        """,
    )

    parser.add_argument(
        "input_file",
        help="Path to the input file (PDF, DOCX, PPTX, or image)",
    )
    parser.add_argument(
        "--output-dir", "-o",
        default="./extraction_output",
        help="Output directory for results (default: ./extraction_output)",
    )
    parser.add_argument(
        "--ocr-lang",
        default="eng",
        help="Tesseract OCR language(s), e.g. 'eng', 'eng+hin' (default: eng)",
    )
    parser.add_argument(
        "--use-ai",
        action="store_true",
        help="Enable AI vision captioning (requires transformers + torch)",
    )
    parser.add_argument(
        "--max-keywords",
        type=int,
        default=15,
        help="Maximum keywords per image (default: 15)",
    )
    parser.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Enable verbose/debug logging",
    )
    parser.add_argument(
        "--check-deps",
        action="store_true",
        help="Check available dependencies and exit",
    )
    parser.add_argument(
        "--json-only",
        action="store_true",
        help="Output only JSON to stdout (no formatting)",
    )
    parser.add_argument(
        "--no-filter",
        action="store_true",
        help="Disable smart content filtering (keep all text and images)",
    )

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    # --- Check dependencies ---
    extractor = ContentExtractor(
        output_dir=args.output_dir,
        ocr_lang=args.ocr_lang,
        use_ai=args.use_ai,
        max_keywords=args.max_keywords,
        smart_filter=not args.no_filter,
    )

    if args.check_deps:
        print("\nDependency Status:")
        print("-" * 40)
        deps = extractor.check_dependencies()
        for name, available in deps.items():
            status = "INSTALLED" if available else "MISSING"
            print(f"  {name:<25} {status}")
        print()
        sys.exit(0)

    # --- Run extraction ---
    try:
        result = extractor.extract(args.input_file)

        if args.json_only:
            print(result.to_json())
        else:
            print_result_summary(result)

    except FileNotFoundError as e:
        logger.error(f"File not found: {e}")
        sys.exit(1)
    except ImportError as e:
        logger.error(f"Missing dependency: {e}")
        sys.exit(1)
    except ValueError as e:
        logger.error(f"Invalid input: {e}")
        sys.exit(1)
    except Exception as e:
        logger.error(f"Extraction failed: {e}", exc_info=args.verbose)
        sys.exit(1)


if __name__ == "__main__":
    main()
