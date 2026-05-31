"""
Document ingestion pipeline for the Academic Knowledge Intelligence Layer.

Pipeline:
  Upload → Format Detection → Text Extraction → Cleaning →
  Semantic Chunking → Academic Classification → Storage

Supported formats: PDF, DOCX, PPTX, TXT, MD, PNG, JPG, JPEG
"""

from __future__ import annotations

import io
import logging
import re
import time
import base64
from pathlib import Path
from typing import Any
from uuid import uuid4

from sqlalchemy.orm import Session

from ..config import settings
from .chunking import AcademicChunk, semantic_chunk, count_tokens
from .classifier import classify_chunk
from ..llm_pipeline import LLMCall
from .models import (
    AcademicDocument,
    ChunkApprovalStatus,
    DocumentType,
    KnowledgeChunk,
    ProcessingStatus,
    SubjectSyllabus,
    ExtractedImage,
)

logger = logging.getLogger("app.academic.ingestion")


# ---------------------------------------------------------------------------
# Text Extraction
# ---------------------------------------------------------------------------

def extract_text_from_pdf(content: bytes) -> tuple[str, int, dict[int, int]]:
    """Extract text from PDF with page tracking.
    
    Returns: (full_text, page_count, char_offset_to_page_map)
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages_text: list[str] = []
    page_offsets: dict[int, int] = {}
    offset = 0

    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        page_offsets[offset] = page_num
        pages_text.append(text)
        offset += len(text) + 1  # +1 for newline

    full_text = "\n".join(pages_text)
    return full_text, len(reader.pages), page_offsets


def extract_text_from_docx(content: bytes) -> tuple[str, int]:
    """Extract text from DOCX files including tables.
    
    Returns: (full_text, estimated_page_count)
    """
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(content))
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)

    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    full_text = "\n".join(parts)
    # Rough page estimate: ~500 words per page
    est_pages = max(1, len(full_text.split()) // 500)
    return full_text, est_pages


def extract_text_from_pptx(content: bytes) -> tuple[str, int]:
    """Extract text from PPTX files.
    
    Returns: (full_text, slide_count)
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, cannot extract PPTX text")
        return "", 0

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text += text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_text += " | ".join(row_texts) + "\n"
        parts.append(slide_text)

    return "\n".join(parts), len(prs.slides)


def extract_text_from_image(content: bytes) -> str:
    """Extract text from images using OCR.
    
    Tries local OCR first, then falls back to the configured vision model.
    """
    try:
        import pytesseract
        from PIL import Image
        image = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(image).strip()
        if text:
            return text
    except ImportError:
        logger.info("pytesseract not installed, image OCR unavailable")
    except Exception as e:
        logger.warning("Image OCR failed: %s", e)

    llm = LLMCall(
        model=settings.ollama_vision_model,
        timeout=settings.ollama_request_timeout_seconds,
    )
    if not llm.is_available():
        return ""

    system = (
        "You are performing OCR for academic material ingestion. "
        "Transcribe the visible text faithfully and do not summarize."
    )
    prompt = (
        "Read the academic image and return only the extracted text. "
        "Preserve headings, lists, and technical terms whenever possible."
    )
    encoded = base64.b64encode(content).decode("utf-8")
    extracted = llm.generate_text(prompt, system, images=[encoded], model=settings.ollama_vision_model)
    return extracted.strip() if extracted else ""


def extract_text(file_name: str, content: bytes) -> tuple[str, int, dict[int, int] | None]:
    """
    Detect file format and extract text.
    
    Returns: (text, page_count, page_offset_map_or_None)
    """
    suffix = Path(file_name).suffix.lower()
    
    if suffix == ".pdf":
        text, pages, offsets = extract_text_from_pdf(content)
        return text, pages, offsets
    elif suffix == ".docx":
        text, pages = extract_text_from_docx(content)
        return text, pages, None
    elif suffix == ".pptx":
        text, slides = extract_text_from_pptx(content)
        return text, slides, None
    elif suffix in {".txt", ".md", ".csv"}:
        text = content.decode("utf-8", errors="ignore")
        pages = max(1, len(text.split()) // 500)
        return text, pages, None
    elif suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
        text = extract_text_from_image(content)
        return text, 1, None
    else:
        text = content.decode("utf-8", errors="ignore")
        return text, 1, None


# ---------------------------------------------------------------------------
# Text Cleaning
# ---------------------------------------------------------------------------

def clean_academic_text(text: str) -> str:
    """Clean extracted text while preserving academic structure."""
    if not text:
        return ""

    # Normalize whitespace within lines (but preserve paragraph breaks)
    lines = text.split("\n")
    cleaned_lines: list[str] = []
    
    for line in lines:
        # Collapse multiple spaces
        line = re.sub(r"[ \t]+", " ", line)
        # Remove control characters except newline
        line = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", line)
        cleaned_lines.append(line.strip())

    result = "\n".join(cleaned_lines)
    # Collapse more than 3 consecutive newlines
    result = re.sub(r"\n{4,}", "\n\n\n", result)
    return result.strip()


# ---------------------------------------------------------------------------
# Document Type Detection
# ---------------------------------------------------------------------------

def detect_document_type(file_name: str, text: str) -> DocumentType:
    """Heuristically detect the type of academic document."""
    name_lower = file_name.lower()
    text_lower = text[:3000].lower() if text else ""

    if any(kw in name_lower for kw in ("syllabus", "curriculum")):
        return DocumentType.SYLLABUS
    if any(kw in name_lower for kw in ("previous", "model paper", "past paper", "question paper")):
        return DocumentType.PREVIOUS_PAPER
    if any(kw in name_lower for kw in ("question bank", "qbank", "q-bank")):
        return DocumentType.QUESTION_BANK
    if any(kw in name_lower for kw in ("lab", "manual", "experiment")):
        return DocumentType.LAB_MANUAL
    if name_lower.endswith(".pptx") or name_lower.endswith(".ppt"):
        return DocumentType.PPT

    # Content-based detection
    if any(kw in text_lower for kw in ("syllabus", "course objectives", "course outcomes")):
        return DocumentType.SYLLABUS
    if re.search(r"(?:Q\.?\s*\d|question\s+paper|marks?:?\s*\d)", text_lower):
        return DocumentType.QUESTION_BANK

    return DocumentType.NOTES


# ---------------------------------------------------------------------------
# Main Ingestion Pipeline
# ---------------------------------------------------------------------------

def create_document_record(
    db: Session,
    subject_id: int,
    user_id: int,
    file_name: str,
    content: bytes,
    document_type: DocumentType | None = None,
) -> AcademicDocument:
    """Save the file and create the initial database record immediately."""
    suffix = Path(file_name).suffix.lower()

    # --- Save file ---
    upload_dir = settings.storage_path / "academic" / str(subject_id)
    upload_dir.mkdir(parents=True, exist_ok=True)
    saved_path = upload_dir / f"{uuid4().hex}_{file_name}"
    saved_path.write_bytes(content)

    # --- Create document record ---
    doc = AcademicDocument(
        subject_id=subject_id,
        uploaded_by=user_id,
        file_name=file_name,
        file_type=suffix.lstrip("."),
        document_type=document_type or DocumentType.NOTES,
        storage_path=str(saved_path),
        processing_status=ProcessingStatus.EXTRACTING,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return doc


def process_document_background(
    document_id: int,
    auto_approve_threshold: float = 0.6,
) -> None:
    """
    Background worker for extraction, chunking, classification, and embedding.
    Runs completely separate from the request thread.
    """
    from ..database import SessionLocal
    db = SessionLocal()
    
    try:
        doc = db.get(AcademicDocument, document_id)
        if not doc:
            return

        start_time = time.time()
        file_name = doc.file_name

        try:
            content = Path(doc.storage_path).read_bytes()

            # --- Advanced Extract text & images via content_extractor ---
            # Set up output directory for extracted images
            subject_images_dir = settings.storage_path / "academic" / str(doc.subject_id)
            subject_images_dir.mkdir(parents=True, exist_ok=True)
            
            from . import content_extractor
            
            # Map suffix to extractor class
            suffix = Path(file_name).suffix.lower()
            temp_output_dir = str(subject_images_dir)
            
            # Temporary local lists to map blocks and page mapping
            text = ""
            page_count = 1
            chunks = []
            extraction_result = None

            if suffix == ".pdf":
                extractor = content_extractor.PDFExtractor(output_dir=temp_output_dir)
            elif suffix == ".docx":
                extractor = content_extractor.DocxExtractor(output_dir=temp_output_dir)
            elif suffix == ".pptx":
                extractor = content_extractor.PPTXExtractor(output_dir=temp_output_dir)
            elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}:
                extractor = content_extractor.ImageExtractor(output_dir=temp_output_dir)
            else:
                extractor = None
                
            if extractor:
                extraction_result = extractor.extract(doc.storage_path)
                # Apply content filter to clean extracted result
                content_filter = content_extractor.ContentFilter(enabled=True)
                extraction_result = content_filter.filter_result(extraction_result)
                
                # Assemble cleaned text
                text_parts = [block.text for block in extraction_result.text_blocks]
                text = "\n\n".join(text_parts)
                page_count = extraction_result.total_pages
            else:
                # Fallback to simple text extraction
                text, page_count, _ = extract_text(file_name, content)
                
            doc.page_count = page_count

            if not text or len(text.strip()) < 50:
                doc.processing_status = ProcessingStatus.FAILED
                doc.processing_error = "Insufficient text extracted from document"
                db.commit()
                return

            # --- Clean text ---
            cleaned = clean_academic_text(text)
            doc.extracted_text = cleaned[:100000]  # Store first 100k chars

            # --- Detect document type ---
            if doc.document_type == DocumentType.NOTES:
                doc.document_type = detect_document_type(file_name, cleaned)

            # --- Semantic chunking ---
            doc.processing_status = ProcessingStatus.CHUNKING
            db.commit()

            # For page number tracking in chunking
            page_offsets = None
            if extraction_result:
                page_offsets = {}
                offset = 0
                for block in extraction_result.text_blocks:
                    page_offsets[offset] = block.page
                    offset += len(block.text) + 2

            chunks = semantic_chunk(
                cleaned,
                min_tokens=180,
                max_tokens=420,
                overlap_ratio=0.1,
                page_numbers=page_offsets,
            )

            if not chunks:
                logger.warning(
                    "Semantic chunking produced no chunks for '%s'; creating a fallback chunk",
                    file_name,
                )
                fallback_text = cleaned[:4000].strip()
                if not fallback_text:
                    doc.processing_status = ProcessingStatus.FAILED
                    doc.processing_error = "No meaningful chunks could be created"
                    db.commit()
                    return
                chunks = [
                    AcademicChunk(
                        text=fallback_text,
                        chunk_index=0,
                        token_count=count_tokens(fallback_text),
                        page_number=1,
                        source_section=None,
                    )
                ]

            # --- Load syllabus for classification ---
            syllabus = db.query(SubjectSyllabus).filter_by(subject_id=doc.subject_id).first()
            syllabus_modules = syllabus.modules_json if syllabus else None

            # --- Classify and store chunks ---
            doc.processing_status = ProcessingStatus.EMBEDDING
            db.commit()

            db_chunks = []
            for chunk in chunks:
                classification = classify_chunk(
                    chunk.text,
                    source_section=chunk.source_section,
                    syllabus_modules=syllabus_modules,
                )

                approval = (
                    ChunkApprovalStatus.AUTO_APPROVED
                    if classification.confidence_score >= auto_approve_threshold
                    else ChunkApprovalStatus.PENDING_REVIEW
                )

                db_chunk = KnowledgeChunk(
                    document_id=doc.id,
                    subject_id=doc.subject_id,
                    chunk_text=chunk.text,
                    chunk_index=chunk.chunk_index,
                    token_count=chunk.token_count,
                    module_number=classification.module_number,
                    syllabus_unit=None,
                    topic_name=classification.topic_name,
                    bloom_level=classification.bloom_level,
                    co_mapping=classification.co_mapping,
                    page_number=chunk.page_number,
                    confidence_score=classification.confidence_score,
                    approval_status=approval,
                )
                db.add(db_chunk)
                db_chunks.append(db_chunk)

            # --- Process Concept Nodes (Stage 1 & 2) ---
            try:
                from .concept_extraction import extract_concept_nodes
                from .models import ConceptNode

                # Chunk the text into ~4000 character segments to avoid LLM context overflow
                MAX_SEGMENT_LEN = 4000
                segments = [cleaned[i:i + MAX_SEGMENT_LEN] for i in range(0, len(cleaned), MAX_SEGMENT_LEN)]
                
                total_nodes = 0
                for segment in segments:
                    extracted_concepts = extract_concept_nodes(segment)
                    for cdata in extracted_concepts:
                        node = ConceptNode(
                            document_id=doc.id,
                            subject_id=doc.subject_id,
                            topic=cdata.get("topic", "Unknown"),
                            module_number=cdata.get("module_number"),
                            node_type=cdata.get("node_type", "main_topic"),
                            difficulty=cdata.get("difficulty", "medium"),
                            content=cdata.get("content", ""),
                            related_topics=cdata.get("related_topics", []),
                            question_patterns=cdata.get("question_patterns", [])
                        )
                        db.add(node)
                        total_nodes += 1
                db.commit()
                logger.info("Extracted %d ConceptNodes from document", total_nodes)
            except Exception as ce:
                logger.error("Failed to extract ConceptNodes: %s", ce)

            # --- Process and save extracted images ---
            if extraction_result and extraction_result.images:
                image_mapper = content_extractor.ImageKeywordMapper()
                for ext_img in extraction_result.images:
                    # Find context text using mapper
                    related_text, context_before, context_after, keywords = image_mapper.map_image_to_context(
                        extraction_result.text_blocks,
                        ext_img.source_page,
                        block_index=-1
                    )
                    
                    # Create ExtractedImage database record
                    db_img = ExtractedImage(
                        document_id=doc.id,
                        subject_id=doc.subject_id,
                        image_path=ext_img.image_path,
                        source_page=ext_img.source_page,
                        width=ext_img.width,
                        height=ext_img.height,
                        keywords=keywords,
                        ai_caption=ext_img.ai_caption or None,
                        image_hash=ext_img.image_hash
                    )
                    db.add(db_img)

            doc.total_chunks = len(chunks)
            db.commit()

            elapsed = time.time() - start_time
            logger.info(
                "Ingested '%s': %d pages, %d chunks, %d images in %.2fs",
                file_name, page_count, len(chunks), len(extraction_result.images) if extraction_result else 0, elapsed,
            )

            # --- Generate Embeddings ---
            from .embeddings import generate_embeddings_batch
            
            texts = [c.chunk_text for c in db_chunks]
            embeddings = generate_embeddings_batch(texts)
            
            for chunk, embedding in zip(db_chunks, embeddings):
                if embedding is not None:
                    chunk.embedding_vector = embedding
                    
            doc.processing_status = ProcessingStatus.COMPLETED
            db.commit()
            logger.info("Completed embedding generation for doc %d synchronously in background thread", doc.id)

        except Exception as e:
            logger.error("Ingestion failed for '%s': %s", file_name, e)
            doc.processing_status = ProcessingStatus.FAILED
            doc.processing_error = str(e)[:2000]
            db.commit()

    finally:
        db.close()


def process_syllabus_background(document_id: int) -> None:
    """
    Specialized pipeline for extracting strict knowledge bounds from a Syllabus.
    """
    from ..database import SessionLocal
    import json

    db = SessionLocal()
    try:
        doc = db.get(AcademicDocument, document_id)
        if not doc:
            return

        start_time = time.time()
        file_name = doc.file_name

        try:
            content = Path(doc.storage_path).read_bytes()
            
            # --- Advanced Extract text via content_extractor ---
            from . import content_extractor
            
            suffix = Path(file_name).suffix.lower()
            temp_output_dir = str(settings.storage_path / "academic" / str(doc.subject_id))
            
            if suffix == ".pdf":
                extractor = content_extractor.PDFExtractor(output_dir=temp_output_dir)
            elif suffix == ".docx":
                extractor = content_extractor.DocxExtractor(output_dir=temp_output_dir)
            elif suffix == ".pptx":
                extractor = content_extractor.PPTXExtractor(output_dir=temp_output_dir)
            elif suffix in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}:
                extractor = content_extractor.ImageExtractor(output_dir=temp_output_dir)
            else:
                extractor = None
                
            if extractor:
                extraction_result = extractor.extract(doc.storage_path)
                text_parts = [block.text for block in extraction_result.text_blocks]
                text = "\n\n".join(text_parts)
                page_count = extraction_result.total_pages
            else:
                text, page_count, _ = extract_text(file_name, content)
                
            doc.page_count = page_count

            if not text or len(text.strip()) < 50:
                raise ValueError("Insufficient text extracted from syllabus")

            cleaned = clean_academic_text(text)
            doc.extracted_text = cleaned[:100000]

            doc.processing_status = ProcessingStatus.CHUNKING # Using this to mean extracting JSON
            db.commit()

            # --- LLM Extraction ---
            llm = LLMCall(
                model=settings.ollama_model,
                timeout=120,
            )
            
            system = (
                "You are an academic extraction engine. Your job is to strictly extract "
                "the syllabus structure from the provided text into a JSON format."
            )
            prompt = f"""
Extract the syllabus structure from the following text into JSON.
Return ONLY valid JSON. No markdown, no commentary.

Expected JSON structure:
{{
  "modules": [
    {{
      "module_number": 1,
      "title": "Introduction",
      "topics": ["Topic 1", "Topic 2"]
    }}
  ],
  "course_outcomes": [
    {{
      "co_code": "CO1",
      "description": "Understand basics"
    }}
  ]
}}

Text:
{cleaned[:15000]}
"""
            result = llm.generate_text(prompt, system)
            
            # Clean up potential markdown formatting from LLM
            if result.startswith("```json"):
                result = result[7:]
            if result.endswith("```"):
                result = result[:-3]
            
            parsed_json = json.loads(result.strip())
            
            syllabus = db.query(SubjectSyllabus).filter_by(subject_id=doc.subject_id).first()
            if not syllabus:
                syllabus = SubjectSyllabus(subject_id=doc.subject_id)
                db.add(syllabus)
                
            syllabus.modules_json = parsed_json.get("modules", [])
            syllabus.co_json = parsed_json.get("course_outcomes", [])
            
            # Mark doc as completed
            doc.processing_status = ProcessingStatus.COMPLETED
            db.commit()

            logger.info(
                "Extracted Syllabus Intelligence for '%s' in %.2fs",
                file_name, time.time() - start_time,
            )

        except Exception as e:
            logger.error("Syllabus extraction failed for '%s': %s", file_name, e)
            doc.processing_status = ProcessingStatus.FAILED
            doc.processing_error = str(e)[:2000]
            db.commit()

    finally:
        db.close()

